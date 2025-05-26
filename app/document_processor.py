from docx import Document
from pptx import Presentation
import PyPDF2

class DocumentProcessor:

    @staticmethod
    def extract_text_from_pdf(file):
        reader = PyPDF2.PdfReader(file)
        return ''.join([page.extract_text() for page in reader.pages])

    @staticmethod
    def extract_text_from_docx(file):
        doc = Document(file)
        return '\n'.join([p.text for p in doc.paragraphs])

    @staticmethod
    def extract_text_from_pptx(file):
        prs = Presentation(file)
        return '\n'.join([
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        ])

    def process(self, uploaded_file):
        mime = uploaded_file.type
        if mime == "application/pdf":
            return self.extract_text_from_pdf(uploaded_file)
        elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return self.extract_text_from_docx(uploaded_file)
        elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self.extract_text_from_pptx(uploaded_file)
        else:
            raise ValueError("Formato de archivo no soportado")
